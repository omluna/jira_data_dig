#!/usr/bin/env python
# -*- coding: utf-8 -*-
import sys
import pprint
import datetime
import pymongo
import pprint
import pandas as pd
import numpy as np
from plotly import tools
import plotly.graph_objs as go
import plotly.offline as py
#import cufflinks as cf
import pptx
from pptx.util import Inches


def make_pptx(file_lists, output_filename):
    pptFile = pptx.Presentation()

    #picFiles = [fn for fn in os.listdir() if fn.endswith('.png')]
    #picFiles = [fn for fn in os.listdir() if fn.endswith('.png')]

    # 按图片编号顺序导入
    for fn in file_lists:
        slide = pptFile.slides.add_slide(pptFile.slide_layouts[6])

        # 为PPTX文件当前幻灯片中第一个文本框设置文字，本文代码中可忽略
        #slide.shapes.placeholders[0].text = fn[:fn.rindex('.')]

        # 导入并为当前幻灯片添加图片，起始位置和尺寸可修改
        slide.shapes.add_picture(fn, Inches(0), Inches(0), Inches(10), Inches(7.5))

    pptFile.save(output_filename)


def bubble_resolved_time(resolved_issues, total_issues, by, dept=None, title='泡泡', height=600, width=1000, bubble_sizeref=2):
    if by == 'assignee':
        mean_times = resolved_issues[resolved_issues['dept'] == dept]['time_fixed'].groupby(resolved_issues['assignee']).mean()
        total_issues_count = total_issues[total_issues['dept'] == dept].groupby(['assignee']).count()['key']
    else:
        mean_times = resolved_issues['time_fixed'].groupby(resolved_issues[by]).mean()
        total_issues_count = total_issues.groupby([by]).count()['key']

    # 有些模块从来没有解决过问题，所以 mean_times 里的 component 跟 total_issues 里的并不相等
    bubble_size = total_issues_count.apply(math.sqrt)
    sizeref = bubble_sizeref*bubble_size.max()/(60**2)
    text_list = total_issues_count[mean_times.index].apply(str).values

    trace0 = go.Scatter(
        x=mean_times.index,
        y=mean_times.values,
        mode='markers+text',
        text=text_list,
        textposition='auto',
        textfont=dict(color='#ffffff'),
        showlegend=False,
        marker=dict(
            color='#005995',
            sizemode='area',
            size=bubble_size[mean_times.index].values,
            sizeref=sizeref,
        )
    )

    mean_time_lt_30 = resolved_issues[resolved_issues['time_fixed'] < 30]['time_fixed'].mean()
    mean_time_total = resolved_issues['time_fixed'].mean()

    layout = go.Layout(
        title=title,
        autosize=False,
        height=height,
        width=width,
        yaxis=dict(
            title='时间（单位:天）',
        ),
        xaxis=dict(tickangle=90),
        margin=dict(b=150),
        annotations=[
            dict(
                x=0,
                y=round(mean_time_lt_30, 2),
                xref='x',
                yref='y',
                text='{}天'.format(round(mean_time_lt_30, 2)),
                font=dict(color='#DB3A34'),
                # showarrow=True,
                # arrowhead=7,
                ax=-50,
                ay=40
            ),
            dict(
                x=0,
                y=round(mean_time_total, 2),
                xref='x',
                yref='y',
                text='{}天'.format(round(mean_time_total, 2)),
                font=dict(color='#B50B1E'),
                # showarrow=True,
                # arrowhead=7,
                ax=-50,
                ay=-50
            )
        ]
    )

    trace1 = go.Scatter(
        x=mean_times.index,
        y=[mean_time_lt_30]*len(mean_times.index),
        mode='lines',
        showlegend=False,
        line=dict(color='#DB3A34')
    )

    trace2 = go.Scatter(
        x=mean_times.index,
        y=[mean_time_total]*len(mean_times.index),
        mode='lines',
        showlegend=False,
        line=dict(color='#B50B1E')
    )

    data = [trace1, trace2, trace0]
    fig = go.Figure(data=data, layout=layout)
    py.plot(fig, image='jpeg', image_width=width, image_height=height)


def bug_count_employee(issues, group, orig_project_list, width=800, row_height=500, filename='count', max_col=3):
    issues_p1 = issues[(issues['priority'] == 'P1-Highest') & (issues['group'] == group)]
    issues_p2 = issues[(issues['priority'] == 'P2-High') & (issues['group'] == group)]

    project_list = orig_project_list[:]
    for project in project_list:
        if len(issues_p1[issues_p1['project'] == project]) + len(issues_p2[issues_p2['project'] == project]) == 0:
            print('no bug for ' + project)
            project_list.remove(project)

    total = len(issues_p1) + len(issues_p2)
    now = datetime.datetime.today().strftime('%Y-%m-%d %H:%M:%S')

    subplot_gap = 1/3
    prj_count = len(project_list)
    print('Generate {} report for {}'.format(group, str(project_list)))
    row, over = divmod(prj_count, max_col)
    xaxis_width = 1/(max_col + (max_col - 1)/3)
    start_row = 1
    bar_list = []
    bar_width = 0.8
    showlegend = True
    if over > 0:
        row = row + 1
        start_row = 2

    yaxis_width = 1/(row + (row - 1)/3)
    height = row_height*row

    layout = go.Layout(
        title='{} BUG<br>合计:{}'.format(group, total),
        autosize=False,
        barmode='stack',
        width=width,
        height=height,
        # bargap=0.4,
        # font=dict(size=25),
        # bargroupgap=0.05,
        annotations=[
                dict(
                    text='auto generate at ' + now,
                    font=dict(color='#b7b6b7'),
                    #opacity=0.9,
                    showarrow=False,
                    x=0.04,
                    y=1.05,
                    xref='paper',
                    yref='paper'
                ),
        ]
    )

    if over > 0:
        for x in range(1, over + 1):
            project = project_list[x - 1]
            issues_pp1 = issues_p1[issues_p1['project'] == project]
            issues_pp2 = issues_p2[issues_p2['project'] == project]
            p1 = issues_pp1.groupby(['assignee']).size().sort_values(ascending=False)
            p2 = issues_pp2.groupby(['assignee']).size().sort_values(ascending=False)

            xaxis = 'x{}'.format(x)
            yaxis = 'y1'
            text_list = list(p1.index)
            for name in p2.index:
                if name not in p1.index:
                    text_list.append(name)

            bar_list.append(go.Bar(x=p1.index, y=p1.values, text=p1.values, width=bar_width, textposition='inside',
                                   name='P1-Highest', textfont=dict(color='#ffffff'), showlegend=showlegend,
                                   marker=dict(color='rgb(239,8,8)',), xaxis=xaxis, yaxis=yaxis))
            bar_list.append(go.Bar(name='P2-High', x=p2.index, y=p2.values, text=p2.values, width=bar_width,
                                   textposition='inside', textfont=dict(color='#ffffff'), showlegend=showlegend,
                                   marker=dict(color='rgb(8,89,156)',), xaxis=xaxis, yaxis=yaxis))
            showlegend = False
            domain = [(x-1)*xaxis_width + (x-1)*subplot_gap*xaxis_width,
                      x*xaxis_width + (x-1)*subplot_gap*xaxis_width]
            update_xaxis = {'xaxis{}'.format(x): {'anchor': 'y1', 'domain': domain, 
                                                    'title': '{} 合计:{}'.format(project, len(issues_pp2) + len(issues_pp1)),
                                                    'ticktext':text_list,
                                                    'tickvals':text_list,
                                                    'tickangle': 90}}
            layout.update(update_xaxis)

        y = 1
        update_yaxis = {'yaxis{}'.format(y): {'domain': [(y-1)*yaxis_width + (y-1)*subplot_gap*yaxis_width,
                                                         y*yaxis_width + (y-1)*subplot_gap*yaxis_width]}, }
        layout.update(update_yaxis)
        project_list = project_list[over:]

    for y in range(start_row, row+1):
        for x in range(1, max_col+1):
            project = project_list[(y-start_row)*max_col+x-1]
            issues_pp1 = issues_p1[issues_p1['project'] == project]
            issues_pp2 = issues_p2[issues_p2['project'] == project]
            p1 = issues_pp1.groupby(['assignee']).size().sort_values(ascending=False)
            p2 = issues_pp2.groupby(['assignee']).size().sort_values(ascending=False)

            xaxis = 'x{}'.format(over+x+(y-start_row)*max_col)
            yaxis = 'y{}'.format(y)
            text_list = list(p1.index)
            for name in p2.index:
                if name not in p1.index:
                    text_list.append(name)


            bar_list.append(go.Bar(x=p1.index, y=p1.values, text=p1.values, textposition='inside', width=bar_width,
                                   name='P1-Highest', textfont=dict(color='#ffffff'), showlegend=showlegend,
                                   marker=dict(color='rgb(239,8,8)',), xaxis=xaxis, yaxis=yaxis))
            bar_list.append(go.Bar(name='P2-High', x=p2.index, y=p2.values, text=p2.values, width=bar_width,
                                   textposition='inside', textfont=dict(color='#ffffff'), showlegend=showlegend,
                                   marker=dict(color='rgb(8,89,156)',), xaxis=xaxis, yaxis=yaxis))
            showlegend = False
            domain = [(x-1)*xaxis_width + (x-1)*subplot_gap*xaxis_width,
                      x*xaxis_width + (x-1)*subplot_gap*xaxis_width]

            update_xaxis = {'xaxis{}'.format(over+x+(y-start_row)*max_col): {'anchor': yaxis,
                                                                             'domain': domain, 
                                                                             'title': '{} 合计:{}'.format(project, len(issues_pp2) + len(issues_pp1)),
                                                                             'ticktext':text_list,
                                                                             'tickvals':text_list,
                                                                             'tickangle': 90}, }
            layout.update(update_xaxis)

        update_yaxis = {'yaxis{}'.format(y): {'domain': [(y-1)*yaxis_width + (y-1)*subplot_gap*yaxis_width,
                                                         y*yaxis_width + (y-1)*subplot_gap*yaxis_width]}, }
        layout.update(update_yaxis)

    fig = go.Figure(data=bar_list, layout=layout)
    # py.iplot(fig)
    py.plot(fig, filename='{}.html'.format(filename), auto_open=False,
            image='png', image_height=height, image_width=width, image_filename=filename)


def bug_trend(issues, title='BUG趋势图', width=1440, height=1080, image_filename='bugtrend'):
    now = datetime.date.today()
    weekday = now.weekday()
    monday = now - datetime.timedelta(weekday)
    resolved_issues = issues[(issues.status.isin(['Resolved', 'Closed'])) & (issues.resolved < monday)]
    create_issue = issues[issues.created < monday].groupby(['created']).count()['key']
    created_by_day = create_issue.resample('W').sum().fillna(0)
    resolved_issued_by_day = resolved_issues.groupby(['resolved']).count()['key'].resample('W').sum().fillna(0)
    unresolved_issue_byday = created_by_day.sub(resolved_issued_by_day, fill_value=0).cumsum()
    
    patch_value = []
    patch_date = []
    if len(created_by_day) == 1:
        patch_date = [(created_by_day.index[0] - 1).strftime('%y-%m-%d')]
        patch_value = [0]

    recent_limit = 0
    if len(resolved_issued_by_day) > 11:
        #just show the recent 12 weeks
        recent_limit = len(resolved_issued_by_day) - 11

    data = [
        go.Scatter(
            y=unresolved_issue_byday.values[recent_limit:],
            x=list(unresolved_issue_byday.index.strftime('%y-%m-%d'))[recent_limit:],
            mode='lines+markers',
            name='未解BUG趋势线',
            # marker=dict(color='#000000'),
        ),

	go.Bar(
            y=(patch_value + list(created_by_day.values))[recent_limit:],
            x=(patch_date + list(created_by_day.index.strftime('%y-%m-%d')))[recent_limit:],
            text=(patch_value + list(created_by_day.values))[recent_limit:],
            textposition='outside',
            name='新提交的BUG',
            marker=dict(
                color='#106a9c',
            ),
            #textfont=dict(size=22),
            outsidetextfont=dict(size=11),
            constraintext='inside'
            #opacity=0.9
        ),

        go.Bar(
            y=unresolved_issue_byday.values[recent_limit:],
            x=list(unresolved_issue_byday.index.strftime('%y-%m-%d'))[recent_limit:],
            text=unresolved_issue_byday.values[recent_limit:],
            textposition='outside',
            name='剩余未解的BUG',
            marker=dict(
                color='#ff9417',
                #color='#07bdff',
            ),
            # textfont=dict(size=22),
            outsidetextfont=dict(size=12),
            constraintext='inside'
            # opacity=0.8
        ),

        go.Bar(
            y=resolved_issued_by_day.values[recent_limit:],
            x=list(resolved_issued_by_day.index.strftime('%y-%m-%d'))[recent_limit:],
            text=resolved_issued_by_day.values[recent_limit:],
            textposition='outside',
            # textfont=dict(size=22),
            outsidetextfont=dict(size=11),
            name='解决的BUG',
            marker=dict(
                color='#3dd9d6',
                #color='#5a94ad',
            ),
            constraintext='inside'
            # opacity=0.9
        ),
    ]

    bandxaxis = go.XAxis(
        range=[(created_by_day.index[0] - 1).strftime('%y-%m-%d'),
            (created_by_day.index[-1] + 1).strftime('%y-%m-%d')][recent_limit:],
        ticks="",
        showticklabels=True,
        ticktext=list(created_by_day.index.strftime('%b %d'))[recent_limit:],
        tickvals=list(created_by_day.index.strftime('%y-%m-%d'))[recent_limit:],
        tickfont=dict(size=16)
    )

    bandyaxis = go.YAxis(
        tickfont=dict(size=16)
    )

    layout = go.Layout(
        title=title,
        font=dict(size=25),
        bargap=0.25,
        bargroupgap=0.05,
        height=height,
        width=width,
        legend=dict(
            x=0.06,
            y=1.1,
            font=dict(size=22)
        ),
        margin=dict(t=180),
        xaxis=bandxaxis,
        yaxis=bandyaxis,
        #         shapes=[
        #         # 1st highlight during Feb 4 - Feb 6
        #             {
        #                 'type': 'rect',
        #                 # x-reference is assigned to the x-values
        #                 'xref': 'x',
        #                 # y-reference is assigned to the plot paper [0,1]
        #                 'yref': 'paper',
        #                 'x0': unresolved_issue_byday.index[0].strftime('%y-%m-%d'),
        #                 'y0': 0,
        #                 'x1': unresolved_issue_byday.index[-2].strftime('%y-%m-%d'),
        #                 'y1': 1,
        #                 'fillcolor': '#d3d3d3',
        #                 'opacity': 0.3,
        #                 'line': {
        #                     'width': 0,
        #                 }
        #             },
        #         ]

    )

    fig = go.Figure(data=data, layout=layout)
    # py.iplot(fig)
    py.plot(fig, filename=image_filename + '.html', auto_open=False,
            image='png', image_height=height, image_width=width, image_filename=image_filename)


def bug_employee_week(issues, dept, image_filename, width=800, height=600,):
    issues_unresovled = issues[issues.status.isin(['Open', 'Reopened','In Progress','Assigned'])]
    issues_resolved = issues[issues.status.isin(['Resolved', 'Closed'])]

    now = datetime.date.today()
    weekday = now.weekday()
    monday = now - datetime.timedelta(weekday) 
    last_sunday = monday - datetime.timedelta(1)
    last_monday = monday - datetime.timedelta(7) 
    unresolved_total = dict(list(issues_unresovled.groupby('dept')))[dept].groupby('assignee').count()['key'].sort_values(ascending=False)
    resolved_issues_lastweek = issues_resolved[(issues_resolved.resolved >= last_monday) & (issues_resolved.resolved < monday) ].sort_values(by='resolved')
    resolved_total = dict(list(resolved_issues_lastweek.groupby('dept')))[dept].groupby('assignee').count()['key'].sort_values(ascending=False)

    title = '{} P1 BUG统计图<br>{}至{}'.format(dept, last_monday.strftime("%m-%d"), 
                                           last_sunday.strftime("%m-%d"),)
    gap = int((max(resolved_total.values)+max(unresolved_total.values))/5)
    ymin = -(int(max(resolved_total.values)/gap) + 1)*gap
    ymax = (int(max(unresolved_total.values)/gap) + 1)*gap
    tick = list(range(ymin, ymax + 1,gap))

    bandyaxis = go.YAxis(
        range=[ymin, ymax],
        ticks="", 
        showticklabels=True,
        ticktext=[abs(i) for i in tick],
        tickvals=tick,
    )

    data = []
    data.append(go.Bar(x=unresolved_total.index, y=unresolved_total.values, 
                       text=unresolved_total.values, marker=dict(color='#0271C6'), 
                       name='未解BUG总数(截止{})'.format(last_sunday.strftime("%m-%d")),
                     ))
    data.append(go.Bar(x=resolved_total.index, y=0 - resolved_total.values, text=resolved_total.values,
                       marker=dict(color='#526372'),name="上周解决BUG总数",opacity=0.5),)

    layout = go.Layout(
        title=title,
        #barmode='stack',
        barmode='overlay',
        bargap=0.5,
        bargroupgap=0.12,
        width=width,
        height=height,
        legend=dict(
            x=0.76,
            y=1.1,
           # font=dict(size=22)
        ),
        margin=dict(pad=8),
        xaxis=dict(ticks='outside',ticklen=10,showgrid=True), 
        yaxis=bandyaxis
    )
    fig = go.Figure(data=data, layout=layout)
    py.plot(fig, filename=image_filename + '.html', auto_open=False,
            image='png', image_height=height, image_width=width, image_filename=image_filename)
    #py.iplot(fig)


def parse_changelog(issue):
    now = datetime.datetime.now()
    anchor_time = []
    
    is_opened_bug = True
    last_assigned = issue['assignee_id']
    change_logs = issue['change_logs']
    last_update_time = issue['created'].to_pydatetime()
    for change_log in change_logs:
        t = dict()
        # change_log 里的datetime 需要加上8小时转换为北京时间
        time_from = change_log['date'] + datetime.timedelta(hours=8)
        #pprint.pprint(time_from)
        for item in change_log['items']:
            if item['field'] == 'assignee':
                t['who'] = item['from']
                t['when'] = time_from
                t['time'] = (time_from - last_update_time).total_seconds()/86400
                t['issue_id'] = issue['key']
                last_assigned = item['to']
                last_update_time = time_from

                anchor_time.append(t)
                break
            if item['to'] == 'Resolved':
                is_opened_bug = False
                t['who'] = change_log['author']
                t['when'] = time_from
                t['issue_id'] = issue['key']
                t['time'] = (time_from - last_update_time).total_seconds()/86400
                last_update_time = time_from

                anchor_time.append(t)
                break

            if item['to'] == 'Closed':
                is_opened_bug = False
                t['who'] = change_log['author']
                t['when'] = time_from
                t['issue_id'] = issue['key']
                t['time'] = (time_from - last_update_time).total_seconds()/86400
                last_update_time = time_from

                anchor_time.append(t)
                break
                
            if item['to'] == 'Reopened':
                t['who'] = change_log['author']
                t['issue_id'] = issue['key']
                t['when'] = time_from
                t['time'] = (time_from - last_update_time).total_seconds()/86400
                last_update_time = time_from
                is_opened_bug = True

                anchor_time.append(t)
                break


    if is_opened_bug:
        t = dict()
        t['who'] = last_assigned
        t['issue_id'] = issue['key']
        t['when'] = now
        t['time'] = (now - last_update_time).total_seconds()/86400
        
        anchor_time.append(t)

    return anchor_time   


def stop_time_count_bubble(name, filename, height=400, width=800):
    days = 14
    now = datetime.date.today()
    em = anchro_time_df[anchro_time_df.displayName == name].sort_values(by='when')
    em = em.set_index('when')
    count_sum = em.resample('D')['time'].count()
    time_sum = em.resample('D')['time'].sum()
    mean_time_single = time_sum/count_sum
    mean_time = time_sum.cumsum()/count_sum.cumsum()
    from_day = now - datetime.timedelta(days=days)
    day_label = pd.date_range(from_day, now)
    
    bubble_size = count_sum.apply(math.sqrt)
    bubble_sizeref=5
    sizeref = bubble_sizeref*bubble_size.max()/(50**2)
    text_list = count_sum[from_day:]

    line_time = go.Scatter(
        x=day_label,
        y=mean_time[from_day:].values,
        mode='markers+text',
        text=text_list,
        textposition='auto',
        textfont=dict(color='#ffffff',size=9),
        showlegend=False,
        marker=dict(
            color='#005995',
            sizemode='area',
            size=bubble_size[from_day:].values,
            sizeref=sizeref,
        )
        
    )
    
    bandxaxis = go.XAxis(
        type='date',
        range=[from_day- datetime.timedelta(days=1), now + datetime.timedelta(days=1)],
        ticks="outside", 
        showticklabels=True,
        ticktext=list(day_label.strftime('%m-%d %a')),
        tickvals=day_label,
        tickfont=dict(size=11),
        #tickfont=dict(size=11,color='#B9B9C3'),
        #ticklen=10,
        tickangle=90,
    )
    
    bandyaxis = go.YAxis(
        #type='date',
        title='累计平均处理时间（单位:天）',
        range=[-0.5, mean_time[from_day:].max() + 1],
        #ticks="outside", 
        #showticklabels=True,
        #ticktext=list(day_label.strftime('%m-%d')),
        #tickvals=day_label,
        tickfont=dict(size=11),
        #tickfont=dict(size=11,color='#B9B9C3'),
        #ticklen=10,
        #tickangle=90
        #color='#FFFFFF'
    )

    data = [line_time]
    layout = go.Layout(
        titlefont=dict(size=15),
        title = name + '<br>近两周处理过的BUG总数：' + str(count_sum[from_day:].sum()),
        xaxis = bandxaxis,
        yaxis = bandyaxis,
        height = height,
        width = width,
        margin = dict(b=120)
    )

    fig = go.Figure(data=data, layout=layout)
    py.plot(fig, filename=filename + '.html', auto_open=False,
            image='png', image_height=height, image_width=width, image_filename=filename)
    #py.iplot(fig)